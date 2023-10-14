"""
Modifed version of the llama-index SlidesReader, which won't attempt to caption images.
"""

import os
from pathlib import Path
from typing import Dict, List, Optional
from pptx import Presentation

from llama_index.readers.base import BaseReader
from llama_index.schema import Document


class PptxReaderNoCaption(BaseReader):
    """Powerpoint parser.

    Extract text, caption images, and specify slides.

    """

    def __init__(self) -> None:
        self.parser_config = {}

    def load_data(
        self,
        file: Path,
        extra_info: Optional[Dict] = None,
    ) -> List[Document]:
        """Parse file."""
        from pptx import Presentation

        presentation = Presentation(file)
        result = ""
        for i, slide in enumerate(presentation.slides):
            result += f"\n\nSlide #{i}: \n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    result += f"{shape.text}\n"

        return [Document(text=result, metadata=extra_info or {})]